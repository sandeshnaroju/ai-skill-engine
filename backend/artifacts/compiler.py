"""
backend/artifacts/compiler.py
Compiles Artifacts into native binary and text files:
- .docx via python-docx (styled tables, headings, code blocks, hyperlinks, images, horizontal dividers)
- .xlsx via openpyxl
- .pptx via python-pptx
- .pdf  via reportlab (styled tables, headings, preformatted code cards, hyperlinks, images, horizontal dividers)
- Plain text / code (.py, .js, .svg, .json, .csv, .html)
"""
import io
import re
import json
import base64
import urllib.request
from typing import Tuple, List, Optional, Any, Dict

from models import SessionArtifact
from .manager import assemble_full_content


def _hex_to_rgb(hex_str: str) -> Tuple[int, int, int]:
    """Helper to convert #hex or #rgb to RGB tuple."""
    if not hex_str:
        return (51, 65, 85)
    hex_clean = hex_str.strip().lstrip('#')
    if len(hex_clean) == 3:
        hex_clean = "".join([c * 2 for c in hex_clean])
    if len(hex_clean) == 6:
        try:
            return (int(hex_clean[0:2], 16), int(hex_clean[2:4], 16), int(hex_clean[4:6], 16))
        except ValueError:
            pass
    return (51, 65, 85)


def _strip_html(text: str) -> str:
    """Strip HTML tags and decode common entities for plain text fallback."""
    if not text:
        return ""
    clean = re.sub(r'<[^<]+?>', '', text)
    clean = clean.replace('&nbsp;', ' ').replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>')
    return clean.strip()


def _fetch_image_bytes(url_or_data: str) -> Optional[bytes]:
    """Fetch image bytes from data URI or HTTP(S) URL with a safe timeout."""
    if not url_or_data:
        return None
    url_clean = url_or_data.strip()
    if url_clean.startswith("data:image/"):
        try:
            if "," in url_clean:
                _, encoded = url_clean.split(",", 1)
                return base64.b64decode(encoded)
        except Exception:
            return None
    elif url_clean.startswith("http://") or url_clean.startswith("https://"):
        try:
            req = urllib.request.Request(
                url_clean,
                headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}
            )
            with urllib.request.urlopen(req, timeout=5) as resp:
                return resp.read()
        except Exception:
            return None
    return None


def _tokenize_inline_formatting(text: str) -> List[dict]:
    """
    Parses inline markdown and HTML tags into tokens.
    Handles:
    - Markdown links [text](url) and <a href="...">text</a>
    - Colored spans <span style="color:..."> and <font color="...">
    - Bold (**bold**, <b>, <strong>)
    - Italic (*italic*, _italic_, <i>, <em>)
    - Inline code (`code`)
    - Plain text
    """
    if not text:
        return []

    # Clean unrenderable layout tags
    clean = re.sub(r'</?(?:div|p|br|mark|section|article)[^>]*>', '', text)

    pattern = re.compile(
        r'(?P<md_link>\[(?P<link_text>[^\]]+)\]\((?P<link_url>[^\)]+)\))|'
        r'(?P<html_link><a\s+[^>]*href=[\'"](?P<a_url>[^\'"]+)[\'"][^>]*>(?P<a_text>.*?)</a>)|'
        r'(?P<color_span><span\s+style=[\'"][^\'"]*color:\s*(?P<span_color>#[0-9a-fA-F]{3,6}|[a-zA-Z]+)[^\'"]*[\'"][^>]*>(?P<span_text>.*?)</span>)|'
        r'(?P<color_font><font\s+color=[\'"](?P<font_color>#[0-9a-fA-F]{3,6}|[a-zA-Z]+)[\'"][^>]*>(?P<font_text>.*?)</font>)|'
        r'(?P<bold_tag><b>|<strong>)(?P<bold_tag_text>.*?)(?:</b>|</strong>)|'
        r'(?P<italic_tag><i>|<em>)(?P<italic_tag_text>.*?)(?:</i>|</em>)|'
        r'(?:\*\*(?P<bold_md>[^\*]+)\*\*)|'
        r'(?:__(?P<bold_under>[^_]+)__)|'
        r'(?:\*(?P<it_md>[^\*]+)\*)|'
        r'(?:_(?P<it_under>[^_]+)_)|'
        r'(?:`(?P<code_txt>[^`]+)`)',
        re.IGNORECASE | re.DOTALL
    )

    tokens = []
    last_idx = 0

    for match in pattern.finditer(clean):
        start, end = match.span()
        if start > last_idx:
            tokens.append({"type": "text", "text": clean[last_idx:start]})

        d = match.groupdict()

        if d.get("link_text") and d.get("link_url"):
            tokens.append({"type": "link", "text": d["link_text"], "url": d["link_url"]})
        elif d.get("a_text") and d.get("a_url"):
            tokens.append({"type": "link", "text": d["a_text"], "url": d["a_url"]})
        elif d.get("span_text") and d.get("span_color"):
            tokens.append({"type": "colored", "color": d["span_color"], "text": d["span_text"]})
        elif d.get("font_text") and d.get("font_color"):
            tokens.append({"type": "colored", "color": d["font_color"], "text": d["font_text"]})
        elif d.get("bold_tag_text"):
            tokens.append({"type": "bold", "text": d["bold_tag_text"]})
        elif d.get("bold_md"):
            tokens.append({"type": "bold", "text": d["bold_md"]})
        elif d.get("bold_under"):
            tokens.append({"type": "bold", "text": d["bold_under"]})
        elif d.get("italic_tag_text"):
            tokens.append({"type": "italic", "text": d["italic_tag_text"]})
        elif d.get("it_md"):
            tokens.append({"type": "italic", "text": d["it_md"]})
        elif d.get("it_under"):
            tokens.append({"type": "italic", "text": d["it_under"]})
        elif d.get("code_txt"):
            tokens.append({"type": "code", "text": d["code_txt"]})

        last_idx = end

    if last_idx < len(clean):
        tokens.append({"type": "text", "text": clean[last_idx:]})

    return tokens


def compile_to_docx(artifact: SessionArtifact) -> io.BytesIO:
    """Compile Markdown/Document artifact into native Microsoft Word (.docx) with real tables, hyperlinks, code blocks, and styled typography."""
    import docx
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml import parse_xml
    from docx.oxml.ns import nsdecls

    doc = docx.Document()

    # Set 1-inch margins
    for section in doc.sections:
        section.top_margin = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin = Inches(1.0)
        section.right_margin = Inches(1.0)

    # Document Title
    title_p = doc.add_paragraph()
    title_run = title_p.add_run(artifact.title)
    title_run.font.name = "Calibri"
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(15, 23, 42)
    title_p.paragraph_format.space_after = Pt(16)

    def add_hyperlink(paragraph, url, text, color="2563EB"):
        """Add a real clickable external hyperlink to a python-docx paragraph."""
        try:
            part = paragraph.part
            r_id = part.relate_to(url, docx.opc.constants.RELATIONSHIP_TYPE.HYPERLINK, is_external=True)
            hyperlink = parse_xml(f'<w:hyperlink {nsdecls("w")} r:id="{r_id}" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"/>')
            new_run = parse_xml(f'<w:r {nsdecls("w")}><w:rPr><w:color w:val="{color}"/><w:u w:val="single"/><w:rFonts w:ascii="Calibri" w:hAnsi="Calibri"/></w:rPr><w:t>{_strip_html(text)}</w:t></w:r>')
            hyperlink.append(new_run)
            paragraph._p.append(hyperlink)
        except Exception:
            run = paragraph.add_run(_strip_html(text))
            run.font.name = "Calibri"
            run.font.color.rgb = RGBColor(37, 99, 235)
            run.font.underline = True

    def apply_inline_tokens_to_paragraph(p, text_content):
        tokens = _tokenize_inline_formatting(text_content)
        if not tokens:
            return
        for t in tokens:
            t_type = t.get("type", "text")
            t_text = _strip_html(t.get("text", ""))
            if not t_text:
                continue

            if t_type == "link":
                add_hyperlink(p, t.get("url", "#"), t_text)
                continue

            run = p.add_run(t_text)
            run.font.name = "Calibri"
            run.font.size = Pt(11)

            if t_type == "bold":
                run.bold = True
                run.font.color.rgb = RGBColor(15, 23, 42)
            elif t_type == "italic":
                run.italic = True
                run.font.color.rgb = RGBColor(51, 65, 85)
            elif t_type == "code":
                run.font.name = "Consolas"
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(124, 58, 237)
            elif t_type == "colored":
                r, g, b = _hex_to_rgb(t.get("color", "#7c3aed"))
                run.font.color.rgb = RGBColor(r, g, b)
                run.bold = True
            else:
                run.font.color.rgb = RGBColor(51, 65, 85)

    def set_cell_shading(cell, color_hex):
        shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color_hex.lstrip("#")}"/>')
        cell._tc.get_or_add_tcPr().append(shading)

    def set_cell_margins(cell, top=100, bottom=100, left=150, right=150):
        tcPr = cell._tc.get_or_add_tcPr()
        tcMar = parse_xml(f'<w:tcMar {nsdecls("w")}><w:top w:w="{top}" w:type="dxa"/><w:bottom w:w="{bottom}" w:type="dxa"/><w:left w:w="{left}" w:type="dxa"/><w:right w:w="{right}" w:type="dxa"/></w:tcMar>')
        tcPr.append(tcMar)

    def add_horizontal_rule():
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(8)
        p.paragraph_format.space_after = Pt(8)
        pPr = p._p.get_or_add_pPr()
        pBdr = parse_xml(f'<w:pBdr {nsdecls("w")}><w:bottom w:val="single" w:sz="6" w:space="1" w:color="CBD5E1"/></w:pBdr>')
        pPr.append(pBdr)

    def add_code_block_table(code_text):
        tbl = doc.add_table(rows=1, cols=1)
        tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
        tblPr = tbl._tbl.tblPr
        borders = parse_xml(f'<w:tblBorders {nsdecls("w")}><w:top w:val="single" w:sz="4" w:space="0" w:color="E2E8F0"/><w:bottom w:val="single" w:sz="4" w:space="0" w:color="E2E8F0"/><w:left w:val="single" w:sz="12" w:space="0" w:color="6366F1"/><w:right w:val="single" w:sz="4" w:space="0" w:color="E2E8F0"/></w:tblBorders>')
        tblPr.append(borders)

        cell = tbl.cell(0, 0)
        set_cell_shading(cell, "F8FAFC")
        set_cell_margins(cell, top=120, bottom=120, left=180, right=180)
        p = cell.paragraphs[0]
        p.paragraph_format.space_after = Pt(0)
        run = p.add_run(code_text)
        run.font.name = "Consolas"
        run.font.size = Pt(9.5)
        run.font.color.rgb = RGBColor(30, 41, 59)

        spacer = doc.add_paragraph()
        spacer.paragraph_format.space_after = Pt(6)

    def add_image_to_document(url_or_data, alt_text=""):
        img_bytes = _fetch_image_bytes(url_or_data)
        if img_bytes:
            try:
                from PIL import Image as PILImage
                pil_img = PILImage.open(io.BytesIO(img_bytes))
                w, h = pil_img.size
                max_width_in = 5.5
                aspect = h / w if w > 0 else 0.75
                img_w_in = min(w / 96.0, max_width_in)
                if img_w_in < 1.0:
                    img_w_in = max_width_in
                img_h_in = img_w_in * aspect
                if img_h_in > 6.0:
                    img_h_in = 6.0
                    img_w_in = img_h_in / aspect

                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(6)
                p.paragraph_format.space_after = Pt(3)
                run = p.add_run()
                run.add_picture(io.BytesIO(img_bytes), width=Inches(img_w_in))

                if alt_text:
                    caption_p = doc.add_paragraph()
                    caption_p.paragraph_format.space_after = Pt(8)
                    caption_run = caption_p.add_run(f"Figure: {_strip_html(alt_text)}")
                    caption_run.font.name = "Calibri"
                    caption_run.font.size = Pt(9)
                    caption_run.font.italic = True
                    caption_run.font.color.rgb = RGBColor(100, 116, 139)
                return
            except Exception:
                pass

        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(6)
        run = p.add_run(f"🖼 [Image: {_strip_html(alt_text or 'Image')}]")
        run.font.name = "Calibri"
        run.font.size = Pt(10)
        run.font.italic = True
        run.font.color.rgb = RGBColor(100, 116, 139)

    blocks = sorted(artifact.blocks, key=lambda b: b.order_index) if artifact.blocks else []
    for block in blocks:
        lines = block.content.splitlines()
        in_code_block = False
        code_block_lines = []
        table_rows = []

        def flush_table():
            nonlocal table_rows
            if not table_rows:
                return
            clean_rows = []
            for row in table_rows:
                if all(re.match(r'^:?-+:?$', c.strip()) for c in row if c.strip()):
                    continue
                clean_rows.append(row)

            if clean_rows:
                cols_count = max(len(r) for r in clean_rows)
                tbl = doc.add_table(rows=len(clean_rows), cols=cols_count)
                tbl.alignment = WD_TABLE_ALIGNMENT.CENTER

                tblPr = tbl._tbl.tblPr
                borders = parse_xml(f'<w:tblBorders {nsdecls("w")}><w:top w:val="single" w:sz="4" w:space="0" w:color="CBD5E1"/><w:bottom w:val="single" w:sz="6" w:space="0" w:color="94A3B8"/><w:left w:val="none"/><w:right w:val="none"/><w:insideH w:val="single" w:sz="4" w:space="0" w:color="E2E8F0"/><w:insideV w:val="none"/></w:tblBorders>')
                tblPr.append(borders)

                for r_idx, row_data in enumerate(clean_rows):
                    is_header = (r_idx == 0)
                    for c_idx in range(cols_count):
                        cell_val = row_data[c_idx] if c_idx < len(row_data) else ""
                        cell = tbl.cell(r_idx, c_idx)
                        set_cell_margins(cell, top=120, bottom=120, left=180, right=180)

                        if is_header:
                            set_cell_shading(cell, "1E293B")
                            p = cell.paragraphs[0]
                            p.alignment = WD_ALIGN_PARAGRAPH.LEFT
                            run = p.add_run(_strip_html(cell_val))
                            run.font.name = "Calibri"
                            run.font.size = Pt(10)
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(255, 255, 255)
                        else:
                            if r_idx % 2 == 1:
                                set_cell_shading(cell, "F8FAFC")
                            else:
                                set_cell_shading(cell, "FFFFFF")
                            p = cell.paragraphs[0]
                            p.alignment = WD_ALIGN_PARAGRAPH.LEFT
                            apply_inline_tokens_to_paragraph(p, cell_val)

                spacer_p = doc.add_paragraph()
                spacer_p.paragraph_format.space_after = Pt(8)

            table_rows = []

        for line in lines:
            raw_line = line
            line_str = line.strip()

            # Handle Code Blocks
            if line_str.startswith("```"):
                if in_code_block:
                    in_code_block = False
                    code_text = "\n".join(code_block_lines)
                    add_code_block_table(code_text)
                    code_block_lines = []
                else:
                    flush_table()
                    in_code_block = True
                    code_block_lines = []
                continue

            if in_code_block:
                code_block_lines.append(raw_line)
                continue

            # Handle Tables
            if line_str.startswith("|") and line_str.endswith("|"):
                cells = [c.strip() for c in line_str.strip("|").split("|")]
                table_rows.append(cells)
                continue
            else:
                flush_table()

            if not line_str:
                continue

            # Horizontal Rule (---, ***, ___, <hr>)
            if re.match(r'^(?:---|___|\*\*\*|\-{3,}|\*{3,}|_{3,}|<hr\s*/?>)$', line_str):
                add_horizontal_rule()
                continue

            # Markdown Image: ![alt](url)
            img_match = re.match(r'^!\[(.*?)\]\((.*?)\)$', line_str)
            if img_match:
                alt_txt, img_url = img_match.group(1), img_match.group(2)
                add_image_to_document(img_url, alt_txt)
                continue

            # HTML Image: <img ... src="..." ...>
            html_img_match = re.search(r'<img\s+[^>]*src=[\'"]([^\'"]+)[\'"][^>]*>', line_str)
            if html_img_match:
                img_url = html_img_match.group(1)
                alt_match = re.search(r'alt=[\'"]([^\'"]+)[\'"]', line_str)
                alt_txt = alt_match.group(1) if alt_match else ""
                add_image_to_document(img_url, alt_txt)
                continue

            # Headings
            if line_str.startswith("# "):
                h = doc.add_paragraph()
                h.paragraph_format.space_before = Pt(16)
                h.paragraph_format.space_after = Pt(6)
                run = h.add_run(_strip_html(line_str[2:].strip()))
                run.font.name = "Calibri"
                run.font.size = Pt(18)
                run.font.bold = True
                run.font.color.rgb = RGBColor(30, 41, 59)
            elif line_str.startswith("## "):
                h = doc.add_paragraph()
                h.paragraph_format.space_before = Pt(12)
                h.paragraph_format.space_after = Pt(4)
                run = h.add_run(_strip_html(line_str[3:].strip()))
                run.font.name = "Calibri"
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(51, 65, 85)
            elif line_str.startswith("### "):
                h = doc.add_paragraph()
                h.paragraph_format.space_before = Pt(10)
                h.paragraph_format.space_after = Pt(3)
                run = h.add_run(_strip_html(line_str[4:].strip()))
                run.font.name = "Calibri"
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = RGBColor(71, 85, 105)
            elif line_str.startswith("- ") or line_str.startswith("* "):
                p = doc.add_paragraph(style='List Bullet')
                p.paragraph_format.space_after = Pt(3)
                apply_inline_tokens_to_paragraph(p, line_str[2:].strip())
            elif re.match(r'^\d+\.\s+', line_str):
                text_val = re.sub(r'^\d+\.\s+', '', line_str)
                p = doc.add_paragraph(style='List Number')
                p.paragraph_format.space_after = Pt(3)
                apply_inline_tokens_to_paragraph(p, text_val)
            elif line_str.startswith(">"):
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = Inches(0.25)
                p.paragraph_format.space_before = Pt(4)
                p.paragraph_format.space_after = Pt(6)
                clean_quote = re.sub(r'^>\s*(\[!NOTE\]|\[!TIP\]|\[!IMPORTANT\])?\s*', '', line_str)
                apply_inline_tokens_to_paragraph(p, clean_quote)
            else:
                p = doc.add_paragraph()
                p.paragraph_format.space_after = Pt(6)
                apply_inline_tokens_to_paragraph(p, line_str)

        flush_table()

    output = io.BytesIO()
    doc.save(output)
    output.seek(0)
    return output


def _parse_spreadsheet_blocks(blocks: List[Any], default_title: str = "Sheet1") -> List[Dict[str, Any]]:
    """
    Parses artifact blocks in any format (JSON, Markdown Table, CSV, TSV) into a list of sheet definitions:
    [{"name": "Sheet 1", "columns": ["Col A", "Col B"], "rows": [["val1", "val2"]]}]
    """
    import csv as pycsv
    sheets = []

    for idx, block in enumerate(blocks):
        raw = str(getattr(block, 'content', '') or '').strip()
        b_title = getattr(block, 'title', '') or f"Sheet {idx + 1}"
        if not raw:
            continue

        # Strip markdown code fences if present (```csv ... ```, ```json ... ```, etc.)
        cleaned = re.sub(r'^```(?:[a-zA-Z0-9_-]+)?\s*\n', '', raw)
        cleaned = re.sub(r'\n```\s*$', '', cleaned).strip()

        # 1. Try parsing JSON format
        try:
            data = json.loads(cleaned)
            if isinstance(data, dict):
                # { "sheets": [ { "sheet_name": "...", "columns": [...], "rows": [...] }, ... ] }
                if "sheets" in data and isinstance(data["sheets"], list) and len(data["sheets"]) > 0:
                    for s_idx, s in enumerate(data["sheets"]):
                        s_name = s.get("sheet_name") or s.get("name") or f"{b_title} {s_idx + 1}"
                        cols = s.get("columns") or []
                        rows = s.get("rows") or []
                        if cols or rows:
                            sheets.append({"name": s_name[:31], "columns": cols, "rows": rows})
                    continue

                # { "sheet_name": "...", "columns": [...], "rows": [...] }
                s_name = data.get("sheet_name") or data.get("name") or b_title
                cols = data.get("columns") or []
                rows = data.get("rows") or []
                if cols or rows:
                    sheets.append({"name": s_name[:31], "columns": cols, "rows": rows})
                    continue

            elif isinstance(data, list) and len(data) > 0:
                # [ { "Col A": "Val 1", "Col B": "Val 2" }, ... ]
                if isinstance(data[0], dict):
                    headers = list(data[0].keys())
                    rows = [[row.get(h, "") for h in headers] for row in data]
                    sheets.append({"name": b_title[:31], "columns": headers, "rows": rows})
                    continue
                # [ ["Col A", "Col B"], ["Val 1", "Val 2"] ]
                elif isinstance(data[0], list):
                    headers = data[0]
                    rows = data[1:]
                    sheets.append({"name": b_title[:31], "columns": headers, "rows": rows})
                    continue
        except Exception:
            pass

        # 2. Try parsing Markdown Table format (| Col 1 | Col 2 |)
        if "|" in cleaned:
            table_lines = [l.strip() for l in cleaned.splitlines() if l.strip().startswith("|") and l.strip().endswith("|")]
            if len(table_lines) >= 2:
                def parse_md_row(line):
                    return [c.strip() for c in line.strip("|").split("|")]

                headers = parse_md_row(table_lines[0])
                data_lines = table_lines[1:]
                # Exclude divider row (e.g. |---|---|)
                if data_lines and all(re.match(r'^:?-+:?$', c) for c in parse_md_row(data_lines[0]) if c):
                    data_lines = data_lines[1:]

                rows = [parse_md_row(l) for l in data_lines]
                sheets.append({"name": b_title[:31], "columns": headers, "rows": rows})
                continue

        # 3. Try CSV / TSV format
        try:
            delimiter = '\t' if '\t' in cleaned else ','
            reader = pycsv.reader(cleaned.splitlines(), delimiter=delimiter)
            all_rows = [r for r in reader if any(field.strip() for field in r)]
            if all_rows:
                headers = [h.strip() for h in all_rows[0]]
                rows = [[c.strip() for c in r] for r in all_rows[1:]]
                sheets.append({"name": b_title[:31], "columns": headers, "rows": rows})
                continue
        except Exception:
            pass

        # 4. Fallback: line-by-line raw text
        lines = [l.strip() for l in cleaned.splitlines() if l.strip()]
        sheets.append({"name": b_title[:31], "columns": ["Content"], "rows": [[l] for l in lines]})

    if not sheets:
        sheets.append({"name": default_title[:31], "columns": ["Column 1", "Column 2"], "rows": [["", ""]]})

    return sheets


def compile_to_xlsx(artifact: SessionArtifact) -> io.BytesIO:
    """Compile spreadsheet into beautifully formatted Microsoft Excel (.xlsx) with styled headers, zebra rows, auto-sized columns, and freeze panes."""
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # Remove initial blank sheet

    blocks = sorted(artifact.blocks, key=lambda b: b.order_index) if artifact.blocks else []
    parsed_sheets = _parse_spreadsheet_blocks(blocks, default_title=artifact.title or "Sheet1")

    # Styling definitions
    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="1E293B", end_color="1E293B", fill_type="solid")  # Dark slate
    data_font = Font(name="Calibri", size=10, color="0F172A")
    alt_fill = PatternFill(start_color="F8FAFC", end_color="F8FAFC", fill_type="solid")  # Subtle zebra
    white_fill = PatternFill(start_color="FFFFFF", end_color="FFFFFF", fill_type="solid")

    thin_border = Border(
        left=Side(style='thin', color='E2E8F0'),
        right=Side(style='thin', color='E2E8F0'),
        top=Side(style='thin', color='E2E8F0'),
        bottom=Side(style='thin', color='E2E8F0')
    )

    for s_data in parsed_sheets:
        ws = wb.create_sheet(title=s_data["name"][:31])
        ws.views.sheetView[0].showGridLines = True

        columns = s_data.get("columns") or []
        rows = s_data.get("rows") or []

        # If columns is empty but rows exist, determine max columns
        if not columns and rows:
            max_c = max(len(r) for r in rows)
            columns = [f"Col {get_column_letter(i+1)}" for i in range(max_c)]

        num_cols = len(columns)
        col_max_lengths = [max(len(str(c)), 8) for c in columns]

        # 1. Render Header Row
        ws.row_dimensions[1].height = 26.0
        for col_idx, col_name in enumerate(columns, 1):
            cell = ws.cell(row=1, column=col_idx, value=_strip_html(str(col_name)))
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            cell.border = thin_border

        # 2. Render Data Rows
        for row_idx, row_values in enumerate(rows, 2):
            ws.row_dimensions[row_idx].height = 20.0
            row_fill = alt_fill if (row_idx % 2 == 1) else white_fill

            for col_idx in range(1, num_cols + 1):
                val = row_values[col_idx - 1] if (col_idx - 1) < len(row_values) else ""
                cell = ws.cell(row=row_idx, column=col_idx)
                cell.font = data_font
                cell.fill = row_fill
                cell.border = thin_border

                # Value formatting & auto-typing
                if val is None or val == "":
                    cell.value = ""
                    cell.alignment = Alignment(vertical="center")
                elif isinstance(val, (int, float)):
                    cell.value = val
                    cell.alignment = Alignment(horizontal="right", vertical="center")
                    if isinstance(val, float):
                        cell.number_format = '#,##0.00'
                else:
                    str_val = _strip_html(str(val)).strip()
                    # Check for integer
                    if re.match(r'^-?\d+$', str_val):
                        try:
                            cell.value = int(str_val)
                            cell.alignment = Alignment(horizontal="right", vertical="center")
                            cell.number_format = '#,##0'
                        except ValueError:
                            cell.value = str_val
                            cell.alignment = Alignment(horizontal="left", vertical="center")
                    # Check for float
                    elif re.match(r'^-?\d+\.\d+$', str_val):
                        try:
                            cell.value = float(str_val)
                            cell.alignment = Alignment(horizontal="right", vertical="center")
                            cell.number_format = '#,##0.00'
                        except ValueError:
                            cell.value = str_val
                            cell.alignment = Alignment(horizontal="left", vertical="center")
                    # Check for Currency (e.g. $1,234.56 or -$50.00)
                    elif re.match(r'^-?\$?\s*([0-9]{1,3}(?:,[0-9]{3})*(?:\.[0-9]+)?|\.[0-9]+)\$?$', str_val):
                        try:
                            clean_num = re.sub(r'[\$,\s]', '', str_val)
                            cell.value = float(clean_num)
                            cell.alignment = Alignment(horizontal="right", vertical="center")
                            cell.number_format = '$#,##0.00'
                        except ValueError:
                            cell.value = str_val
                            cell.alignment = Alignment(horizontal="left", vertical="center")
                    # Check for Percentage (e.g. 25.5%)
                    elif re.match(r'^-?\d+(?:\.\d+)?%$', str_val):
                        try:
                            clean_pct = str_val.rstrip('%')
                            cell.value = float(clean_pct) / 100.0
                            cell.alignment = Alignment(horizontal="right", vertical="center")
                            cell.number_format = '0.0%'
                        except ValueError:
                            cell.value = str_val
                            cell.alignment = Alignment(horizontal="left", vertical="center")
                    else:
                        cell.value = str_val
                        cell.alignment = Alignment(horizontal="left", vertical="center")

                # Track column width
                str_len = len(str(cell.value or ''))
                if (col_idx - 1) < len(col_max_lengths):
                    col_max_lengths[col_idx - 1] = max(col_max_lengths[col_idx - 1], str_len)

        # 3. Freeze top header row so it stays visible during scrolling
        ws.freeze_panes = "A2"

        # 4. Auto-fit column widths
        for col_idx, max_l in enumerate(col_max_lengths, 1):
            col_letter = get_column_letter(col_idx)
            ws.column_dimensions[col_letter].width = min(max(max_l + 4, 12), 48)

    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    return output


def _parse_presentation_blocks(blocks: List[Any], default_title: str = "Presentation") -> List[Dict[str, Any]]:
    """
    Parses slide blocks from JSON, Markdown, or HTML into structured slide dictionaries:
    [{"title": "...", "subtitle": "...", "bullets": [...], "stats": [...], "cards": [...], "steps": [...], "quote": "...", "layout": "...", "notes": "..."}]
    """
    slides = []

    for idx, block in enumerate(blocks):
        raw = str(getattr(block, 'content', '') or '').strip()
        b_title = getattr(block, 'title', '') or f"Slide {idx + 1}"
        if not raw:
            continue

        cleaned = re.sub(r'^```(?:[a-zA-Z0-9_-]+)?\s*\n', '', raw)
        cleaned = re.sub(r'\n```\s*$', '', cleaned).strip()

        # 1. Try parsing JSON format
        try:
            data = json.loads(cleaned)
            if isinstance(data, dict):
                # { "slides": [ { ... }, ... ] }
                if "slides" in data and isinstance(data["slides"], list) and len(data["slides"]) > 0:
                    for s_idx, s in enumerate(data["slides"]):
                        s_title = s.get("title") or f"{b_title} - Slide {s_idx + 1}"
                        slides.append({
                            "title": s_title,
                            "subtitle": s.get("subtitle") or s.get("tagline") or "",
                            "bullets": s.get("bullets") or (s.get("content") if isinstance(s.get("content"), list) else ([s.get("content")] if s.get("content") else [])),
                            "stats": s.get("stats") or [],
                            "cards": s.get("cards") or [],
                            "steps": s.get("steps") or s.get("timeline") or [],
                            "columns": s.get("columns") or [],
                            "quote": s.get("quote") or "",
                            "author": s.get("author") or "",
                            "role": s.get("role") or "",
                            "notes": s.get("notes") or s.get("speaker_notes") or "",
                            "badge": s.get("badge") or s.get("tag") or (f"KEYNOTE" if idx == 0 and s_idx == 0 else ""),
                            "layout": s.get("layout") or ""
                        })
                    continue

                # Single slide JSON
                slides.append({
                    "title": data.get("title") or b_title,
                    "subtitle": data.get("subtitle") or data.get("tagline") or "",
                    "bullets": data.get("bullets") or (data.get("content") if isinstance(data.get("content"), list) else ([data.get("content")] if data.get("content") else [])),
                    "stats": data.get("stats") or [],
                    "cards": data.get("cards") or [],
                    "steps": data.get("steps") or data.get("timeline") or [],
                    "columns": data.get("columns") or [],
                    "quote": data.get("quote") or "",
                    "author": data.get("author") or "",
                    "role": data.get("role") or "",
                    "notes": data.get("notes") or data.get("speaker_notes") or "",
                    "badge": data.get("badge") or data.get("tag") or (f"KEYNOTE" if idx == 0 else ""),
                    "layout": data.get("layout") or ""
                })
                continue
        except Exception:
            pass

        # 2. Parse Markdown Slide format
        lines = [l.strip() for l in cleaned.splitlines() if l.strip()]
        title = b_title
        subtitle = ""
        bullets = []
        quote = ""
        stats = []

        for line in lines:
            if line.startswith("# "):
                title = line[2:].strip()
            elif line.startswith("## ") and not subtitle:
                subtitle = line[3:].strip()
            elif line.startswith("- ") or line.startswith("* ") or line.startswith("• "):
                bullet_txt = line[2:].strip()
                # Check for stat format: **$500k**: Revenue Growth
                stat_m = re.match(r'^\*\*(.+?)\*\*:\s*(.+)$', bullet_txt)
                if stat_m:
                    stats.append({"value": stat_m.group(1), "label": stat_m.group(2)})
                else:
                    bullets.append(bullet_txt)
            elif re.match(r'^\d+\.\s+', line):
                bullets.append(re.sub(r'^\d+\.\s+', '', line))
            elif line.startswith(">"):
                quote = re.sub(r'^>\s*', '', line).strip()
            elif not subtitle:
                subtitle = line
            else:
                bullets.append(line)

        slides.append({
            "title": title,
            "subtitle": subtitle,
            "bullets": bullets,
            "stats": stats,
            "cards": [],
            "steps": [],
            "columns": [],
            "quote": quote,
            "author": "",
            "role": "",
            "notes": "",
            "badge": "KEYNOTE" if idx == 0 else "",
            "layout": ""
        })

    if not slides:
        slides.append({"title": default_title, "subtitle": "AI Generated Presentation", "bullets": [], "stats": [], "cards": [], "steps": [], "columns": [], "quote": "", "author": "", "role": "", "notes": "", "badge": "KEYNOTE", "layout": "hero"})

    return slides


def _add_clean_multiline_text(
    text_frame,
    text,
    font_name="Calibri",
    font_size=None,
    font_bold=False,
    font_italic=False,
    font_color=None,
    space_before=None,
    space_after=None,
    align=None,
    prefix=""
):
    """
    Safely writes multiline text into DrawingML text frames.
    DrawingML strict schema requires single-line text runs inside paragraph (<a:p>) elements.
    Raw newlines inside <a:t> cause Apple Keynote import failure ('file format is invalid').
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    if text is None:
        return
    text_str = str(text).replace("\r\n", "\n").replace("\r", "\n")
    raw_lines = [_strip_html(l).strip() for l in text_str.split("\n")]
    clean_lines = [l for l in raw_lines if l]
    if not clean_lines:
        return

    for idx, line in enumerate(clean_lines):
        if idx == 0 and len(text_frame.paragraphs) == 1 and not text_frame.paragraphs[0].text:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        if space_before is not None and idx == 0:
            p.space_before = space_before
        if space_after is not None and idx == len(clean_lines) - 1:
            p.space_after = space_after
        if align is not None:
            p.alignment = align

        clean_run_text = f"{prefix}{line}".replace("\n", " ").replace("\r", "").strip()
        run = p.add_run()
        run.text = clean_run_text
        if font_name:
            run.font.name = font_name
        if font_size:
            run.font.size = font_size
        run.font.bold = font_bold
        run.font.italic = font_italic
        if font_color:
            run.font.color.rgb = font_color


def compile_to_pptx(artifact: SessionArtifact) -> io.BytesIO:
    """Compile Presentation artifact into 100% compliant 16:9 Microsoft PowerPoint (.pptx) compatible with Apple Keynote, Google Slides and MS Office."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    # 16:9 widescreen layout (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    blocks = sorted(artifact.blocks, key=lambda b: b.order_index) if artifact.blocks else []
    slides_data = _parse_presentation_blocks(blocks, default_title=artifact.title or "Presentation Deck")

    # Theme colors
    bg_color = RGBColor(11, 15, 25)          # #0B0F19 Dark Navy Canvas
    card_bg = RGBColor(22, 31, 48)           # #161F30 Dark Slate Container
    card_border = RGBColor(51, 65, 85)       # #334155 Card Border
    accent_indigo = RGBColor(129, 140, 248)  # #818CF8 Accent Indigo
    accent_cyan = RGBColor(56, 189, 248)     # #38BDF8 Accent Sky
    text_white = RGBColor(255, 255, 255)     # #FFFFFF Main Text
    text_slate = RGBColor(203, 213, 225)     # #CBD5E1 Subtitle / Body
    text_muted = RGBColor(148, 163, 184)     # #94A3B8 Muted Text

    for idx, slide_info in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)

        # 1. Native Slide Background (100% Keynote compliant)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        title_text = _strip_html(slide_info.get("title", f"Slide {idx + 1}"))
        subtitle_text = _strip_html(slide_info.get("subtitle", ""))
        bullets = slide_info.get("bullets", [])
        stats = slide_info.get("stats", [])
        cards = slide_info.get("cards", [])
        steps = slide_info.get("steps", [])
        columns = slide_info.get("columns", [])
        quote_text = _strip_html(slide_info.get("quote", ""))
        badge_text = _strip_html(slide_info.get("badge", ""))
        layout = slide_info.get("layout", "").lower()

        # Determine layout
        if not layout:
            if idx == 0 and not stats and not cards and len(bullets) <= 1:
                layout = "hero"
            elif stats or any(re.search(r'\b(\$?\d+[\d,\.]*[%kMGBx\+]*)\b', str(b)) for b in bullets):
                layout = "stats"
            elif quote_text:
                layout = "quote"
            elif steps:
                layout = "timeline"
            elif cards or len(bullets) >= 4:
                layout = "grid"
            elif columns or len(bullets) in (2, 3):
                layout = "split"
            else:
                layout = "standard"

        # ── HERO / TITLE SLIDE ──
        if layout in ("hero", "title", "title_slide"):
            # Accent decorative bar
            accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(1.8), Inches(0.8), Inches(0.08))
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = accent_indigo
            accent_bar.line.color.rgb = accent_indigo
            accent_bar.line.width = Pt(1)

            # Category Pill Badge
            if badge_text:
                badge_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(2.15), Inches(2.2), Inches(0.45))
                badge_shape.fill.solid()
                badge_shape.fill.fore_color.rgb = card_bg
                badge_shape.line.color.rgb = accent_indigo
                badge_shape.line.width = Pt(1)
                b_tf = badge_shape.text_frame
                b_tf.word_wrap = True
                _add_clean_multiline_text(
                    b_tf,
                    badge_text.upper(),
                    font_size=Pt(11),
                    font_bold=True,
                    font_color=accent_cyan,
                    align=PP_ALIGN.CENTER
                )

            # Main Title & Subtitle Box
            top_pos = 2.8 if badge_text else 2.2
            tx_box = slide.shapes.add_textbox(Inches(1.2), Inches(top_pos), Inches(10.9), Inches(3.8))
            tf = tx_box.text_frame
            tf.word_wrap = True

            _add_clean_multiline_text(
                tf,
                title_text,
                font_size=Pt(38),
                font_bold=True,
                font_color=text_white,
                space_after=Pt(14)
            )

            if subtitle_text:
                _add_clean_multiline_text(
                    tf,
                    subtitle_text,
                    font_size=Pt(18),
                    font_color=text_slate
                )

        # ── CONTENT SLIDES (STATS, CARDS, TIMELINE, QUOTE, SPLIT, STANDARD) ──
        else:
            # Header Section
            hdr_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.3))
            htf = hdr_box.text_frame
            htf.word_wrap = True

            _add_clean_multiline_text(
                htf,
                title_text,
                font_size=Pt(28),
                font_bold=True,
                font_color=text_white
            )

            if subtitle_text:
                _add_clean_multiline_text(
                    htf,
                    subtitle_text,
                    font_size=Pt(14),
                    font_color=text_muted,
                    space_before=Pt(4)
                )

            # ── STATS LAYOUT ──
            if layout == "stats" and (stats or bullets):
                items = stats if stats else [{"value": str(b).split(":", 1)[0], "label": str(b).split(":", 1)[1] if ":" in str(b) else "Key Metric"} for b in bullets]
                num_items = min(len(items), 4)
                card_w = (11.333 - (0.4 * (num_items - 1))) / num_items

                for s_idx, st in enumerate(items[:num_items]):
                    left_pos = 1.0 + s_idx * (card_w + 0.4)
                    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(2.2), Inches(card_w), Inches(4.2))
                    card.fill.solid()
                    card.fill.fore_color.rgb = card_bg
                    card.line.color.rgb = card_border
                    card.line.width = Pt(1)

                    ctf = card.text_frame
                    ctf.word_wrap = True

                    _add_clean_multiline_text(
                        ctf,
                        str(st.get("value", "")),
                        font_size=Pt(34),
                        font_bold=True,
                        font_color=accent_indigo,
                        space_after=Pt(8)
                    )
                    _add_clean_multiline_text(
                        ctf,
                        str(st.get("label", "")),
                        font_size=Pt(15),
                        font_bold=True,
                        font_color=text_white
                    )
                    if st.get("desc"):
                        _add_clean_multiline_text(
                            ctf,
                            str(st["desc"]),
                            font_size=Pt(12),
                            font_color=text_muted,
                            space_before=Pt(6)
                        )

            # ── CARDS / GRID LAYOUT ──
            elif layout == "grid" and (cards or bullets):
                items = cards if cards else [{"title": f"Feature {b_i+1}", "desc": str(b)} for b_i, b in enumerate(bullets)]
                num_items = min(len(items), 4)
                card_w = (11.333 - (0.4 * (num_items - 1))) / num_items

                for c_idx, cd in enumerate(items[:num_items]):
                    left_pos = 1.0 + c_idx * (card_w + 0.4)
                    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(2.2), Inches(card_w), Inches(4.2))
                    card.fill.solid()
                    card.fill.fore_color.rgb = card_bg
                    card.line.color.rgb = card_border
                    card.line.width = Pt(1)

                    ctf = card.text_frame
                    ctf.word_wrap = True

                    _add_clean_multiline_text(
                        ctf,
                        str(cd.get("title", f"Point {c_idx+1}")),
                        font_size=Pt(17),
                        font_bold=True,
                        font_color=accent_cyan,
                        space_after=Pt(8)
                    )
                    _add_clean_multiline_text(
                        ctf,
                        str(cd.get("desc", cd.get("text", ""))),
                        font_size=Pt(13),
                        font_color=text_slate
                    )

            # ── TIMELINE / STEPS LAYOUT ──
            elif layout == "timeline" and (steps or bullets):
                items = steps if steps else [{"title": f"Step {s_i+1}", "desc": str(b)} for s_i, b in enumerate(bullets)]
                num_items = min(len(items), 4)
                step_w = (11.333 - (0.4 * (num_items - 1))) / num_items

                for st_idx, st in enumerate(items[:num_items]):
                    left_pos = 1.0 + st_idx * (step_w + 0.4)
                    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(2.2), Inches(step_w), Inches(4.2))
                    card.fill.solid()
                    card.fill.fore_color.rgb = card_bg
                    card.line.color.rgb = accent_indigo
                    card.line.width = Pt(1.5)

                    ctf = card.text_frame
                    ctf.word_wrap = True

                    _add_clean_multiline_text(
                        ctf,
                        f"PHASE {st_idx + 1}",
                        font_size=Pt(11),
                        font_bold=True,
                        font_color=accent_cyan,
                        space_after=Pt(6)
                    )
                    _add_clean_multiline_text(
                        ctf,
                        str(st.get("title", f"Step {st_idx + 1}")),
                        font_size=Pt(16),
                        font_bold=True,
                        font_color=text_white,
                        space_after=Pt(6)
                    )
                    _add_clean_multiline_text(
                        ctf,
                        str(st.get("desc", "")),
                        font_size=Pt(12.5),
                        font_color=text_slate
                    )

            # ── QUOTE LAYOUT ──
            elif layout == "quote" and (quote_text or bullets):
                q_content = quote_text if quote_text else (bullets[0] if bullets else "")
                q_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.2), Inches(10.333), Inches(4.0))
                q_card.fill.solid()
                q_card.fill.fore_color.rgb = card_bg
                q_card.line.color.rgb = accent_indigo
                q_card.line.width = Pt(2)

                qtf = q_card.text_frame
                qtf.word_wrap = True

                _add_clean_multiline_text(
                    qtf,
                    f'"{q_content}"',
                    font_size=Pt(22),
                    font_italic=True,
                    font_color=text_white,
                    space_after=Pt(14)
                )

                author = _strip_html(slide_info.get("author", ""))
                role = _strip_html(slide_info.get("role", ""))
                if author or role:
                    author_str = f"— {author}" + (f", {role}" if role else "")
                    _add_clean_multiline_text(
                        qtf,
                        author_str,
                        font_size=Pt(14),
                        font_bold=True,
                        font_color=accent_indigo
                    )

            # ── SPLIT / COLUMNS LAYOUT ──
            elif layout == "split" and (columns or len(bullets) >= 2):
                col_items = columns if columns else [
                    {"title": "Key Points", "bullets": bullets[:len(bullets)//2 + len(bullets)%2]},
                    {"title": "Key Takeaways", "bullets": bullets[len(bullets)//2 + len(bullets)%2:]}
                ]
                col_w = (11.333 - 0.5) / 2

                for col_idx, col_data in enumerate(col_items[:2]):
                    left_pos = 1.0 + col_idx * (col_w + 0.5)
                    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(2.2), Inches(col_w), Inches(4.2))
                    card.fill.solid()
                    card.fill.fore_color.rgb = card_bg
                    card.line.color.rgb = card_border
                    card.line.width = Pt(1)

                    ctf = card.text_frame
                    ctf.word_wrap = True

                    _add_clean_multiline_text(
                        ctf,
                        str(col_data.get("title", f"Column {col_idx + 1}")),
                        font_size=Pt(17),
                        font_bold=True,
                        font_color=accent_indigo,
                        space_after=Pt(10)
                    )

                    for b_str in col_data.get("bullets", []):
                        _add_clean_multiline_text(
                            ctf,
                            str(b_str),
                            font_size=Pt(13.5),
                            font_color=text_slate,
                            space_after=Pt(6),
                            prefix="• "
                        )

            # ── STANDARD BULLET LIST ──
            else:
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.2))
                card.fill.solid()
                card.fill.fore_color.rgb = card_bg
                card.line.color.rgb = card_border
                card.line.width = Pt(1)

                ctf = card.text_frame
                ctf.word_wrap = True

                if bullets:
                    for b_item in bullets:
                        _add_clean_multiline_text(
                            ctf,
                            str(b_item),
                            font_size=Pt(15),
                            font_color=text_slate,
                            space_after=Pt(8),
                            prefix="• "
                        )
                else:
                    _add_clean_multiline_text(
                        ctf,
                        str(slide_info.get("content", "")),
                        font_size=Pt(15),
                        font_color=text_slate
                    )

        # ── SPEAKER NOTES ──
        notes = slide_info.get("notes", "")
        if notes:
            try:
                notes_slide = slide.notes_slide
                notes_tf = notes_slide.notes_text_frame
                _add_clean_multiline_text(
                    notes_tf,
                    str(notes),
                    font_size=Pt(12),
                    font_color=RGBColor(0, 0, 0)
                )
            except Exception:
                pass

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def compile_to_pdf(artifact: SessionArtifact) -> io.BytesIO:
    """Compile Markdown/Document into high-quality PDF with native ReportLab tables, code blocks, hyperlinks, images, and styled typography."""
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table as RLTable, TableStyle, Preformatted, HRFlowable, Image as RLImage
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors

    output = io.BytesIO()
    doc = SimpleDocTemplate(output, pagesize=letter, rightMargin=48, leftMargin=48, topMargin=48, bottomMargin=48)
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=22,
        leading=26,
        textColor=colors.HexColor("#0F172A"),
        spaceAfter=14
    )
    h1_style = ParagraphStyle(
        'DocH1',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=15,
        leading=19,
        textColor=colors.HexColor("#1E293B"),
        spaceBefore=12,
        spaceAfter=5
    )
    h2_style = ParagraphStyle(
        'DocH2',
        parent=styles['Heading3'],
        fontName='Helvetica-Bold',
        fontSize=12,
        leading=16,
        textColor=colors.HexColor("#334155"),
        spaceBefore=8,
        spaceAfter=4
    )
    body_style = ParagraphStyle(
        'DocBody',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9.5,
        leading=14,
        textColor=colors.HexColor("#334155"),
        spaceAfter=6
    )
    bullet_style = ParagraphStyle(
        'DocBullet',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9.5,
        leading=14,
        leftIndent=15,
        textColor=colors.HexColor("#334155"),
        spaceAfter=3
    )
    tbl_hdr_style = ParagraphStyle(
        'TblHdr',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=8.5,
        leading=11,
        textColor=colors.HexColor("#FFFFFF")
    )
    tbl_cell_style = ParagraphStyle(
        'TblCell',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=8.5,
        leading=11,
        textColor=colors.HexColor("#1E293B")
    )
    code_style = ParagraphStyle(
        'DocCodeStyle',
        parent=styles['Code'],
        fontName='Courier',
        fontSize=8.5,
        leading=11,
        textColor=colors.HexColor("#0F172A")
    )
    caption_style = ParagraphStyle(
        'ImgCaption',
        parent=styles['Normal'],
        fontName='Helvetica-Oblique',
        fontSize=8,
        textColor=colors.HexColor("#64748B"),
        spaceAfter=6,
        spaceBefore=2
    )

    def convert_markdown_inline_to_reportlab_xml(text_content):
        """Convert markdown, links, and HTML color spans into ReportLab formatted XML tags."""
        if not text_content:
            return ""
        tokens = _tokenize_inline_formatting(text_content)
        xml_runs = []
        for t in tokens:
            t_type = t.get("type", "text")
            t_txt = _strip_html(t.get("text", ""))
            if not t_txt:
                continue
            safe = t_txt.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

            if t_type == "link":
                url = t.get("url", "#").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                xml_runs.append(f"<font color='#2563eb'><u><a href='{url}'>{safe}</a></u></font>")
            elif t_type == "bold":
                xml_runs.append(f"<b>{safe}</b>")
            elif t_type == "italic":
                xml_runs.append(f"<i>{safe}</i>")
            elif t_type == "code":
                xml_runs.append(f"<font color='#7c3aed' face='Courier'>{safe}</font>")
            elif t_type == "colored":
                c = t.get("color", "#7c3aed")
                xml_runs.append(f"<font color='{c}'><b>{safe}</b></font>")
            else:
                xml_runs.append(safe)
        return "".join(xml_runs)

    def add_pdf_image(url_or_data, alt_text=""):
        img_bytes = _fetch_image_bytes(url_or_data)
        if img_bytes:
            try:
                from PIL import Image as PILImage
                pil_img = PILImage.open(io.BytesIO(img_bytes))
                w, h = pil_img.size
                max_w = 480.0
                aspect = h / w if w > 0 else 0.75
                img_w = min(w, max_w)
                img_h = img_w * aspect
                if img_h > 450:
                    img_h = 450
                    img_w = img_h / aspect
                story.append(Spacer(1, 4))
                story.append(RLImage(io.BytesIO(img_bytes), width=img_w, height=img_h))
                if alt_text:
                    story.append(Paragraph(f"Figure: {_strip_html(alt_text)}", caption_style))
                else:
                    story.append(Spacer(1, 6))
                return
            except Exception:
                pass

        story.append(Paragraph(f"[Image: {_strip_html(alt_text or 'Image')}]", caption_style))

    story = [Paragraph(artifact.title, title_style), Spacer(1, 8)]

    blocks = sorted(artifact.blocks, key=lambda b: b.order_index) if artifact.blocks else []
    for block in blocks:
        lines = block.content.splitlines()
        in_code_block = False
        code_block_lines = []
        table_rows = []

        def flush_pdf_table():
            nonlocal table_rows
            if not table_rows:
                return
            clean_rows = []
            for row in table_rows:
                if all(re.match(r'^:?-+:?$', c.strip()) for c in row if c.strip()):
                    continue
                clean_rows.append(row)

            if clean_rows:
                cols_count = max(len(r) for r in clean_rows)
                flowable_matrix = []

                for r_idx, row_data in enumerate(clean_rows):
                    is_header = (r_idx == 0)
                    row_cells = []
                    for c_idx in range(cols_count):
                        cell_val = row_data[c_idx] if c_idx < len(row_data) else ""
                        if is_header:
                            row_cells.append(Paragraph(_strip_html(cell_val), tbl_hdr_style))
                        else:
                            row_cells.append(Paragraph(convert_markdown_inline_to_reportlab_xml(cell_val), tbl_cell_style))
                    flowable_matrix.append(row_cells)

                col_width = 516.0 / cols_count
                rl_tbl = RLTable(flowable_matrix, colWidths=[col_width] * cols_count)
                rl_tbl.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor("#1E293B")),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.HexColor("#FFFFFF")),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                    ('TOPPADDING', (0, 0), (-1, -1), 5),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
                    ('LEFTPADDING', (0, 0), (-1, -1), 6),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor("#E2E8F0")),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.HexColor("#FFFFFF"), colors.HexColor("#F8FAFC")]),
                ]))
                story.append(rl_tbl)
                story.append(Spacer(1, 8))

            table_rows = []

        for line in lines:
            raw_line = line
            line_str = line.strip()

            if line_str.startswith("```"):
                if in_code_block:
                    in_code_block = False
                    code_text = "\n".join(code_block_lines)
                    p_code = Preformatted(code_text, code_style)
                    code_box = RLTable([[p_code]], colWidths=[516])
                    code_box.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor("#F8FAFC")),
                        ('BOX', (0, 0), (-1, -1), 0.5, colors.HexColor("#E2E8F0")),
                        ('LINELEFT', (0, 0), (-1, -1), 2.5, colors.HexColor("#6366F1")),
                        ('TOPPADDING', (0, 0), (-1, -1), 6),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                        ('LEFTPADDING', (0, 0), (-1, -1), 8),
                        ('RIGHTPADDING', (0, 0), (-1, -1), 8),
                    ]))
                    story.append(code_box)
                    story.append(Spacer(1, 6))
                    code_block_lines = []
                else:
                    flush_pdf_table()
                    in_code_block = True
                    code_block_lines = []
                continue

            if in_code_block:
                code_block_lines.append(raw_line)
                continue

            if line_str.startswith("|") and line_str.endswith("|"):
                cells = [c.strip() for c in line_str.strip("|").split("|")]
                table_rows.append(cells)
                continue
            else:
                flush_pdf_table()

            if not line_str:
                continue

            # Horizontal Rule (---, ***, ___, <hr>)
            if re.match(r'^(?:---|___|\*\*\*|\-{3,}|\*{3,}|_{3,}|<hr\s*/?>)$', line_str):
                story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#CBD5E1"), spaceBefore=6, spaceAfter=8))
                continue

            # Markdown Image: ![alt](url)
            img_match = re.match(r'^!\[(.*?)\]\((.*?)\)$', line_str)
            if img_match:
                alt_txt, img_url = img_match.group(1), img_match.group(2)
                add_pdf_image(img_url, alt_txt)
                continue

            # HTML Image: <img ... src="..." ...>
            html_img_match = re.search(r'<img\s+[^>]*src=[\'"]([^\'"]+)[\'"][^>]*>', line_str)
            if html_img_match:
                img_url = html_img_match.group(1)
                alt_match = re.search(r'alt=[\'"]([^\'"]+)[\'"]', line_str)
                alt_txt = alt_match.group(1) if alt_match else ""
                add_pdf_image(img_url, alt_txt)
                continue

            if line_str.startswith("# "):
                story.append(Paragraph(_strip_html(line_str[2:].strip()), h1_style))
            elif line_str.startswith("## "):
                story.append(Paragraph(_strip_html(line_str[3:].strip()), h2_style))
            elif line_str.startswith("### "):
                story.append(Paragraph(_strip_html(line_str[4:].strip()), h2_style))
            elif line_str.startswith("- ") or line_str.startswith("* "):
                bullet_xml = f"&bull; {convert_markdown_inline_to_reportlab_xml(line_str[2:].strip())}"
                story.append(Paragraph(bullet_xml, bullet_style))
            elif re.match(r'^\d+\.\s+', line_str):
                num_prefix = re.match(r'^(\d+\.)\s+', line_str).group(1)
                text_val = re.sub(r'^\d+\.\s+', '', line_str)
                num_xml = f"{num_prefix} {convert_markdown_inline_to_reportlab_xml(text_val)}"
                story.append(Paragraph(num_xml, bullet_style))
            elif line_str.startswith(">"):
                clean_quote = re.sub(r'^>\s*(\[!NOTE\]|\[!TIP\]|\[!IMPORTANT\])?\s*', '', line_str)
                quote_xml = f"<i>{convert_markdown_inline_to_reportlab_xml(clean_quote)}</i>"
                quote_box = RLTable([[Paragraph(quote_xml, body_style)]], colWidths=[516])
                quote_box.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor("#F8FAFC")),
                    ('LINELEFT', (0, 0), (-1, -1), 2.5, colors.HexColor("#6366F1")),
                    ('TOPPADDING', (0, 0), (-1, -1), 4),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
                    ('LEFTPADDING', (0, 0), (-1, -1), 8),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 8),
                ]))
                story.append(quote_box)
                story.append(Spacer(1, 4))
            else:
                xml_p = convert_markdown_inline_to_reportlab_xml(line_str)
                story.append(Paragraph(xml_p, body_style))

        flush_pdf_table()

    doc.build(story)
    output.seek(0)
    return output


def export_artifact(artifact: SessionArtifact, target_format: str = None) -> Tuple[bytes, str, str]:
    """
    Export artifact in requested format.
    Returns: (file_bytes, mime_type, filename)
    """
    ext = target_format or (artifact.filename.split(".")[-1] if "." in artifact.filename else "txt")
    ext = ext.lower().lstrip(".")
    base_name = artifact.filename.rsplit(".", 1)[0] if "." in artifact.filename else artifact.filename

    if ext == "docx" or (ext == "doc" and artifact.artifact_type == "document"):
        buf = compile_to_docx(artifact)
        return buf.read(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document", f"{base_name}.docx"

    elif ext in ("xlsx", "xls") or (ext == "csv" and artifact.artifact_type == "spreadsheet"):
        if ext == "csv":
            content = assemble_full_content(artifact)
            return content.encode("utf-8"), "text/csv", f"{base_name}.csv"
        buf = compile_to_xlsx(artifact)
        return buf.read(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", f"{base_name}.xlsx"

    elif ext == "pptx" or artifact.artifact_type == "presentation":
        buf = compile_to_pptx(artifact)
        return buf.read(), "application/vnd.openxmlformats-officedocument.presentationml.presentation", f"{base_name}.pptx"

    elif ext == "pdf":
        buf = compile_to_pdf(artifact)
        return buf.read(), "application/pdf", f"{base_name}.pdf"

    # Default: raw text/code
    full_text = assemble_full_content(artifact)
    mime = "text/plain"
    if ext == "py":
        mime = "text/x-python"
    elif ext in ("js", "ts"):
        mime = "application/javascript"
    elif ext == "json":
        mime = "application/json"
    elif ext == "svg":
        mime = "image/svg+xml"
    elif ext == "html":
        mime = "text/html"

    out_name = f"{base_name}.{ext}" if not artifact.filename.endswith(f".{ext}") else artifact.filename
    return full_text.encode("utf-8"), mime, out_name
